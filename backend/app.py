from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from .summarizer import summarize_text

app = FastAPI(title="AI Lecture Summarizer")

# Allow frontend requests
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/health")
def health():
    return {"status": "ok"}

@app.post("/api/summarize/pptx")
async def summarize_pptx(file: UploadFile = File(...)):
    if not file.filename.endswith(".pptx"):
        return {"error": "Please upload a .pptx file"}

    prs = Presentation(file.file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)

    full_text = "\n".join(text_runs).strip()
    if not full_text:
        return {"error": "No text found in PowerPoint"}

    summary = summarize_text(full_text)
    return {"summary": summary, "slides": len(prs.slides)}
